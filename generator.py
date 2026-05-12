"""
Proposal generator: takes client info + selected locations.

Output PPT structure:
  Slide 1: KHUSHI cover (from OOH_2/khushi_wrapper.pptx slide 1)
  Slide 2: Client info billboard (from OOH_2 slide 2, SAMPLE TEXT replaced)
  Slides 3+: GSRTC location slides (2 per selected location, from template.pptx)
  Final slide: KHUSHI THANK YOU (from OOH_2 slide 4)

Output XLSX: rate card filtered to selected locations + client header + totals.
"""

import re
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import List

from lxml import etree
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


LOCATION_SLIDE_MAP = {
    "ANAND":         (4, 5),
    "KHEDA":         (6, 7),
    "GANDHINAGAR":   (8, 9),
    "HIMMATNAGAR":   (10, 11),
    "DAHOD":         (12, 13),
    "GODHRA":        (14, 15),
    "ANKLESHWAR":    (16, 17),
    "SURAT":         (18, 19),
    "VALSAD":        (20, 21),
    "CHIKHLI":       (22, 23),
    "BOTAD":         (24, 25),
    "BHAVNAGAR":     (26, 27),
    "VERAVAL":       (28, 29),
    "JUNAGADH":      (30, 31),
    "PORBANDAR":     (32, 33),
    "DWARKA":        (34, 35),
    "JAMNAGAR":      (36, 37),
    "MORBI":         (38, 39),
    "BHACHAU":       (40, 41),
    "SURENDRANAGAR": (42, 43),
}

LOCATIONS_IN_ORDER = list(LOCATION_SLIDE_MAP.keys())
LOCATION_EXCEL_ROW = {loc: i + 2 for i, loc in enumerate(LOCATIONS_IN_ORDER)}

ASSETS_DIR = Path(__file__).parent / "assets"
KHUSHI_WRAPPER = ASSETS_DIR / "khushi_wrapper.pptx"
GSRTC_TEMPLATE = ASSETS_DIR / "template.pptx"
RATECARD = ASSETS_DIR / "ratecard.xlsx"


def sanitize_filename(name: str) -> str:
    name = re.sub(r'[^\w\s-]', '', name).strip()
    name = re.sub(r'[-\s]+', '_', name)
    return name or "Client"


# ============================================================================
# PPT generation
# ============================================================================

def _copy_slide_into_presentation(src_prs: Presentation, src_idx: int, dst_prs: Presentation):
    """
    Append a deep-copy of src_prs.slides[src_idx] onto dst_prs.
    Copies shapes, slide background, and pulls in referenced media + layout.
    """
    src_slide = src_prs.slides[src_idx]
    src_layout = src_slide.slide_layout

    # Find a usable blank layout in dst (we won't try to copy the layout itself —
    # the slide content is self-contained in shapes).
    blank_layout = dst_prs.slide_layouts[6] if len(dst_prs.slide_layouts) > 6 else dst_prs.slide_layouts[0]
    new_slide = dst_prs.slides.add_slide(blank_layout)

    # Strip any placeholder shapes from the new slide
    spTree = new_slide.shapes._spTree
    SHAPE_TAGS = {'sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'}
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag in SHAPE_TAGS:
            spTree.remove(child)

    # Deep-copy each shape from source
    for child in src_slide.shapes._spTree:
        tag = etree.QName(child).localname
        if tag in SHAPE_TAGS:
            spTree.append(deepcopy(child))

    # Copy slide background if present
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    src_cSld = src_slide.element.find(f'{{{P_NS}}}cSld')
    new_cSld = new_slide.element.find(f'{{{P_NS}}}cSld')
    if src_cSld is not None and new_cSld is not None:
        src_bg = src_cSld.find(f'{{{P_NS}}}bg')
        if src_bg is not None:
            existing_bg = new_cSld.find(f'{{{P_NS}}}bg')
            if existing_bg is not None:
                new_cSld.remove(existing_bg)
            new_cSld.insert(0, deepcopy(src_bg))

    # Copy needed relationships + media
    _copy_relationships(src_slide, new_slide)

    return new_slide


def _copy_relationships(src_slide, new_slide):
    """
    Walk new_slide's XML looking for r:embed / r:link attributes, look up the
    source rel, copy media into dst package, and rewrite rIds to new values.
    """
    src_rels = src_slide.part.rels
    new_part = new_slide.part

    rId_map = {}

    for el in new_slide.element.iter():
        for attr_name in list(el.attrib.keys()):
            if not (attr_name.endswith('}embed') or attr_name.endswith('}link') or attr_name.endswith('}id')):
                continue
            # The r:id attribute is also used for things like hyperlinks
            # We only care about attributes in the relationships namespace
            if 'relationships' not in attr_name:
                continue
            old_rId = el.attrib[attr_name]
            if not old_rId or not old_rId.startswith('rId'):
                continue
            if old_rId in rId_map:
                el.attrib[attr_name] = rId_map[old_rId]
                continue
            if old_rId not in src_rels:
                # Not all r:id attrs map to slide rels (some are layout refs we don't carry)
                continue
            src_rel = src_rels[old_rId]
            try:
                if src_rel.is_external:
                    new_rId = new_part.relate_to(
                        src_rel.target_ref, src_rel.reltype, is_external=True
                    )
                else:
                    target = src_rel.target_part
                    new_rId = new_part.relate_to(target, src_rel.reltype)
                rId_map[old_rId] = new_rId
                el.attrib[attr_name] = new_rId
            except Exception:
                # If we can't copy a particular rel, leave it (will render blank)
                pass


def _replace_sample_text_with_client_info(prs: Presentation, client_name: str,
                                          company: str, mobile: str, email: str,
                                          sender_name: str = "", sender_mobile: str = "",
                                          sender_email: str = ""):
    """
    On slide 2 of the KHUSHI template, replace 'SAMPLE TEXT' with a 'Prepared for:'
    block listing company, contact person, mobile, email, followed by a
    'Prepared by:' block with the sender's details.
    """
    if len(prs.slides) < 2:
        return
    slide = prs.slides[1]

    target_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text or ""
        if "SAMPLE TEXT" in text:
            target_shape = shape
            break

    if target_shape is None:
        return

    # Resize the text frame to accommodate the multi-line content.
    # Original is a thin strip (~0.5" tall) sized for one line of "SAMPLE TEXT".
    # We need ~4.5" of vertical space for both blocks at the new font sizes.
    from pptx.util import Inches
    target_shape.top = Inches(1.4)
    target_shape.height = Inches(5.3)
    # Keep left/width as-is so it stays in the white area next to the billboard

    tf = target_shape.text_frame

    # Capture font name from existing first run if possible
    orig_font_name = "Aptos"
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.name:
                orig_font_name = run.font.name
                break
        if orig_font_name != "Aptos":
            break

    # Build lines: (text, size_pt, bold, color_rgb_tuple)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GREY = RGBColor(0x33, 0x33, 0x33)
    MED_GREY = RGBColor(0x55, 0x55, 0x55)
    BRAND_RED = RGBColor(0xC0, 0x00, 0x00)

    lines = []
    # Prepared for block
    lines.append(("Prepared for:", 22, True, BRAND_RED))
    if company.strip():
        lines.append((company.strip(), 26, True, BLACK))
    if client_name.strip():
        lines.append((client_name.strip(), 18, False, DARK_GREY))
    if mobile.strip():
        lines.append((f"Mobile: {mobile.strip()}", 15, False, MED_GREY))
    if email.strip():
        lines.append((f"Email: {email.strip()}", 15, False, MED_GREY))

    # Spacer line
    has_sender = any([sender_name.strip(), sender_mobile.strip(), sender_email.strip()])
    if has_sender:
        lines.append(("", 10, False, MED_GREY))  # spacer
        lines.append(("Prepared by:", 20, True, BRAND_RED))
        if sender_name.strip():
            lines.append((sender_name.strip(), 18, True, BLACK))
        if sender_mobile.strip():
            lines.append((f"Mobile: {sender_mobile.strip()}", 15, False, MED_GREY))
        if sender_email.strip():
            lines.append((f"Email: {sender_email.strip()}", 15, False, MED_GREY))

    # Clear & rebuild
    tf.clear()
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = orig_font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _delete_slide_by_index(prs: Presentation, index: int):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if index < 0 or index >= len(slides):
        return
    slide_id = slides[index]
    rId = slide_id.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    xml_slides.remove(slide_id)
    prs.part.drop_rel(rId)


def _move_slide_to_end(prs: Presentation, current_index: int):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if current_index >= len(slides):
        return
    sld = slides[current_index]
    xml_slides.remove(sld)
    xml_slides.append(sld)


def _renumber_slide_parts(prs: Presentation):
    """
    Rename slide parts so their filenames are slide1.xml..slideN.xml matching
    the order they appear in the presentation. This avoids collisions when
    python-pptx adds new slides after we've deleted/moved some.
    """
    from pptx.opc.packuri import PackURI

    # Collect slides in presentation order via their rIds
    sldIdLst = prs.slides._sldIdLst
    R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

    # Get the (rId -> part) mapping for slide rels on the presentation part
    pres_part = prs.part
    slide_parts_in_order = []
    for sldId in sldIdLst:
        rId = sldId.attrib[R_NS]
        slide_parts_in_order.append(pres_part.rels[rId].target_part)

    # Two-phase rename: first to temp names, then to final, to avoid collisions
    package = pres_part.package
    parts_dict = package._parts if hasattr(package, '_parts') else None

    # Strategy: use the internal package API to rename
    # We change part.partname (which has a setter)
    temp_names = []
    for i, part in enumerate(slide_parts_in_order):
        new_uri = PackURI(f"/ppt/slides/__tmp_{i+1}__.xml")
        temp_names.append(new_uri)
        part.partname = new_uri

    for i, part in enumerate(slide_parts_in_order):
        new_uri = PackURI(f"/ppt/slides/slide{i+1}.xml")
        part.partname = new_uri


def generate_pptx(client_name: str, company: str, mobile: str, email: str,
                  selected_locations: List[str],
                  sender_name: str = "", sender_mobile: str = "",
                  sender_email: str = "") -> bytes:
    selected_set = {loc.upper() for loc in selected_locations}
    invalid = selected_set - set(LOCATION_SLIDE_MAP.keys())
    if invalid:
        raise ValueError(f"Unknown locations: {invalid}")

    khushi = Presentation(str(KHUSHI_WRAPPER))
    gsrtc = Presentation(str(GSRTC_TEMPLATE))

    # KHUSHI starts with [0]=cover, [1]=billboard w/ SAMPLE TEXT, [2]=blank, [3]=thank-you
    # Step 1: Replace SAMPLE TEXT on slide 2
    _replace_sample_text_with_client_info(
        khushi, client_name, company, mobile, email,
        sender_name, sender_mobile, sender_email,
    )

    # Step 2: Delete the blank slide 3 (index 2). Thank-you slide now at index 2.
    _delete_slide_by_index(khushi, 2)

    # Renumber so the existing slide parts are slide1.xml, slide2.xml, slide3.xml
    # (avoids collision when add_slide picks the next name).
    _renumber_slide_parts(khushi)

    # Step 3: Append location slides from GSRTC in template order
    for loc in LOCATIONS_IN_ORDER:
        if loc not in selected_set:
            continue
        photo_idx, stats_idx = LOCATION_SLIDE_MAP[loc]
        _copy_slide_into_presentation(gsrtc, photo_idx, khushi)
        _copy_slide_into_presentation(gsrtc, stats_idx, khushi)

    # Step 4: Move the thank-you slide (still at index 2) to the very end
    _move_slide_to_end(khushi, current_index=2)

    # Final renumber for cleanliness
    _renumber_slide_parts(khushi)

    buf = BytesIO()
    khushi.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================================
# Excel generation
# ============================================================================

def generate_xlsx(client_name: str, company: str, mobile: str, email: str,
                  selected_locations: List[str],
                  sender_name: str = "", sender_mobile: str = "",
                  sender_email: str = "") -> bytes:
    selected_set = {loc.upper() for loc in selected_locations}

    wb = load_workbook(str(RATECARD))
    ws = wb["Sheet1"]
    header_row = [c.value for c in ws[1]]

    selected_rows_data = []
    for loc in LOCATIONS_IN_ORDER:
        if loc in selected_set:
            row_idx = LOCATION_EXCEL_ROW[loc]
            row_data = []
            for c in range(1, len(header_row) + 1):
                row_data.append(ws.cell(row=row_idx, column=c).value)
            selected_rows_data.append((row_idx, row_data))

    new_wb = Workbook()
    new_ws = new_wb.active
    new_ws.title = "Proposal"

    thin = Side(border_style="thin", color="888888")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    red_fill = PatternFill("solid", fgColor="C00000")
    light_fill = PatternFill("solid", fgColor="FCE4E4")
    sender_fill = PatternFill("solid", fgColor="F0F0F0")
    grey_fill = PatternFill("solid", fgColor="D9D9D9")
    white_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")

    new_ws["A1"] = "GSRTC LED SCREEN PROPOSAL"
    new_ws["A1"].font = Font(name="Calibri", size=18, bold=True, color="FFFFFF")
    new_ws["A1"].fill = red_fill
    new_ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    new_ws.merge_cells("A1:M1")
    new_ws.row_dimensions[1].height = 36

    # --- Prepared For block ---
    # Use B for label, C-G for value (left side), H for label, I-M for value (right side)
    new_ws["A2"] = "PREPARED FOR"
    new_ws["A2"].font = Font(name="Calibri", size=13, bold=True, color="C00000")
    new_ws["A2"].fill = light_fill
    new_ws["A2"].alignment = Alignment(horizontal="left", vertical="center", indent=1)
    new_ws.merge_cells("A2:M2")
    new_ws.row_dimensions[2].height = 24

    # Row 3: Company | Contact Person
    new_ws["A3"] = "Company:";        new_ws["C3"] = company
    new_ws["H3"] = "Contact Person:"; new_ws["J3"] = client_name
    # Row 4: Mobile | Email
    new_ws["A4"] = "Mobile:";         new_ws["C4"] = mobile
    new_ws["H4"] = "Email:";          new_ws["J4"] = email

    for cell in ["A3", "H3", "A4", "H4"]:
        new_ws[cell].font = Font(name="Calibri", size=12, bold=True)
        new_ws[cell].fill = light_fill
        new_ws[cell].alignment = Alignment(horizontal="left", vertical="center", indent=1)
    for cell in ["C3", "J3", "C4", "J4"]:
        new_ws[cell].font = Font(name="Calibri", size=12)
        new_ws[cell].alignment = Alignment(horizontal="left", vertical="center")

    new_ws.merge_cells("A3:B3"); new_ws.merge_cells("C3:G3")
    new_ws.merge_cells("H3:I3"); new_ws.merge_cells("J3:M3")
    new_ws.merge_cells("A4:B4"); new_ws.merge_cells("C4:G4")
    new_ws.merge_cells("H4:I4"); new_ws.merge_cells("J4:M4")
    new_ws.row_dimensions[3].height = 22
    new_ws.row_dimensions[4].height = 22

    # --- Prepared By block ---
    has_sender = any([sender_name.strip(), sender_mobile.strip(), sender_email.strip()])
    if has_sender:
        new_ws["A5"] = "PREPARED BY"
        new_ws["A5"].font = Font(name="Calibri", size=13, bold=True, color="C00000")
        new_ws["A5"].fill = sender_fill
        new_ws["A5"].alignment = Alignment(horizontal="left", vertical="center", indent=1)
        new_ws.merge_cells("A5:M5")
        new_ws.row_dimensions[5].height = 24

        new_ws["A6"] = "Name:";   new_ws["C6"] = sender_name
        new_ws["H6"] = "Mobile:"; new_ws["J6"] = sender_mobile
        new_ws["A7"] = "Email:";  new_ws["C7"] = sender_email

        for cell in ["A6", "H6", "A7"]:
            new_ws[cell].font = Font(name="Calibri", size=12, bold=True)
            new_ws[cell].fill = sender_fill
            new_ws[cell].alignment = Alignment(horizontal="left", vertical="center", indent=1)
        for cell in ["C6", "J6", "C7"]:
            new_ws[cell].font = Font(name="Calibri", size=12)
            new_ws[cell].alignment = Alignment(horizontal="left", vertical="center")

        new_ws.merge_cells("A6:B6"); new_ws.merge_cells("C6:G6")
        new_ws.merge_cells("H6:I6"); new_ws.merge_cells("J6:M6")
        new_ws.merge_cells("A7:B7"); new_ws.merge_cells("C7:M7")
        new_ws.row_dimensions[6].height = 22
        new_ws.row_dimensions[7].height = 22
        header_row_num = 9
    else:
        header_row_num = 6

    for col_idx, header in enumerate(header_row, start=1):
        cell = new_ws.cell(row=header_row_num, column=col_idx, value=header)
        cell.font = white_font
        cell.fill = red_fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
    new_ws.row_dimensions[header_row_num].height = 40

    current_row = header_row_num + 1
    for original_row_idx, row_data in selected_rows_data:
        for col_idx, value in enumerate(row_data, start=1):
            if isinstance(value, str) and value.startswith("="):
                value = re.sub(
                    rf'([A-Z]+){original_row_idx}\b',
                    rf'\g<1>{current_row}',
                    value
                )
            cell = new_ws.cell(row=current_row, column=col_idx, value=value)
            cell.font = Font(name="Calibri", size=11)
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border
        current_row += 1

    total_row = current_row
    new_ws.cell(row=total_row, column=2, value="TOTAL")
    new_ws.cell(row=total_row, column=2).font = Font(name="Calibri", size=12, bold=True)
    new_ws.cell(row=total_row, column=2).fill = grey_fill
    new_ws.cell(row=total_row, column=2).alignment = Alignment(horizontal="center")
    new_ws.cell(row=total_row, column=2).border = border

    first_data_row = header_row_num + 1
    last_data_row = total_row - 1
    sum_columns = [5, 6, 10, 12, 13]
    for col_idx in range(1, 14):
        cell = new_ws.cell(row=total_row, column=col_idx)
        cell.fill = grey_fill
        cell.border = border
        cell.alignment = Alignment(horizontal="center")
        if col_idx in sum_columns and last_data_row >= first_data_row:
            col_letter = cell.column_letter
            cell.value = f"=SUM({col_letter}{first_data_row}:{col_letter}{last_data_row})"
            cell.font = Font(name="Calibri", size=12, bold=True)

    widths = {"A": 6, "B": 22, "C": 10, "D": 8, "E": 16, "F": 14, "G": 12,
              "H": 12, "I": 14, "J": 14, "K": 10, "L": 16, "M": 18}
    for col, w in widths.items():
        new_ws.column_dimensions[col].width = w

    # Page setup so the file prints/PDFs cleanly fitting one page wide (landscape)
    new_ws.page_setup.orientation = new_ws.ORIENTATION_LANDSCAPE
    new_ws.page_setup.fitToWidth = 1
    new_ws.page_setup.fitToHeight = 0
    new_ws.sheet_properties.pageSetUpPr.fitToPage = True
    new_ws.print_options.horizontalCentered = True
    new_ws.page_margins.left = 0.3
    new_ws.page_margins.right = 0.3
    new_ws.page_margins.top = 0.4
    new_ws.page_margins.bottom = 0.4

    # --- Terms & Conditions block ---
    tc_start_row = total_row + 2

    new_ws.cell(row=tc_start_row, column=1, value="TERMS & CONDITIONS")
    new_ws.cell(row=tc_start_row, column=1).font = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
    new_ws.cell(row=tc_start_row, column=1).fill = red_fill
    new_ws.cell(row=tc_start_row, column=1).alignment = Alignment(horizontal="center", vertical="center")
    new_ws.merge_cells(start_row=tc_start_row, start_column=1, end_row=tc_start_row, end_column=13)
    new_ws.row_dimensions[tc_start_row].height = 26

    tc_items = [
        "Ads will be displayed on LED screens for a minimum of 30 days. Each ad plays 100 times per day per screen (10 seconds per slot, 5-minute loop).",
        "Reports and audience information will be shared. Screen breakdown is subject to technical issues; alternate screens may be used if required.",
        "Ad creatives must be shared at least 7 days before campaign start and must comply with all legal guidelines (no religious or objectionable content).",
        "Payment Terms: 50% advance to confirm booking; balance within 15 days of campaign start. All prices are in INR + 18% GST. Late payment attracts 1.5% interest per month. Refunds apply only if issues exceed 20% downtime.",
        "All rights, tools, and reports remain with us. Ads may be used for promotional purposes.",
        "Campaign cancellation requires 15 days' prior notice. Early termination is non-refundable.",
        "Our liability is limited to the amount paid. We are not responsible for indirect or additional losses.",
        "All terms are governed by Indian laws.",
    ]

    tc_thin = Side(border_style="thin", color="CCCCCC")
    tc_border = Border(left=tc_thin, right=tc_thin, top=tc_thin, bottom=tc_thin)

    for i, item in enumerate(tc_items, start=1):
        row = tc_start_row + i
        num_cell = new_ws.cell(row=row, column=1, value=i)
        num_cell.font = Font(name="Calibri", size=11, bold=True, color="C00000")
        num_cell.alignment = Alignment(horizontal="center", vertical="top")
        num_cell.fill = PatternFill("solid", fgColor="FFF8F8")
        num_cell.border = tc_border

        text_cell = new_ws.cell(row=row, column=2, value=item)
        text_cell.font = Font(name="Calibri", size=11)
        text_cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
        text_cell.fill = PatternFill("solid", fgColor="FFF8F8")
        text_cell.border = tc_border

        # Merge columns B through M for the text content
        new_ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=13)

        # Set row height proportional to text length (rough estimate)
        # ~85 chars per line at the merged width
        char_count = len(item)
        lines_needed = max(1, (char_count // 85) + (1 if char_count % 85 else 0))
        new_ws.row_dimensions[row].height = max(22, lines_needed * 18)

        # Apply border to the merged cells too (openpyxl needs explicit border on each)
        for col in range(2, 14):
            new_ws.cell(row=row, column=col).border = tc_border

    buf = BytesIO()
    new_wb.save(buf)
    buf.seek(0)
    return buf.read()


if __name__ == "__main__":
    test_locations = ["ANAND", "GANDHINAGAR", "SURAT"]
    print("Generating test PPT...")
    pptx_bytes = generate_pptx(
        client_name="Anil Sharma",
        company="Sharma Industries Pvt Ltd",
        mobile="9876543210",
        email="anil@sharma-industries.com",
        selected_locations=test_locations,
        sender_name="Priya Patel",
        sender_mobile="9988776655",
        sender_email="priya@khushiooh.com",
    )
    Path("/tmp/test_out.pptx").write_bytes(pptx_bytes)
    print(f"  PPT size: {len(pptx_bytes):,} bytes")

    print("Generating test Excel...")
    xlsx_bytes = generate_xlsx(
        client_name="Anil Sharma",
        company="Sharma Industries Pvt Ltd",
        mobile="9876543210",
        email="anil@sharma-industries.com",
        selected_locations=test_locations,
        sender_name="Priya Patel",
        sender_mobile="9988776655",
        sender_email="priya@khushiooh.com",
    )
    Path("/tmp/test_out.xlsx").write_bytes(xlsx_bytes)
    print(f"  Excel size: {len(xlsx_bytes):,} bytes")
    print("Done.")
